from __future__ import annotations

from datetime import date
from io import BytesIO
from typing import Any

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Pt


def format_won(value: float | int) -> str:
    return f"{float(value or 0):,.0f}원"


def set_cell_shading(cell: Any, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = tc_pr.find(qn("w:shd"))
    if shd is None:
        shd = OxmlElement("w:shd")
        tc_pr.append(shd)
    shd.set(qn("w:fill"), fill)


def set_korean_font(document: Document, font_name: str = "맑은 고딕") -> None:
    style = document.styles["Normal"]
    style.font.name = font_name
    style.font.size = Pt(10)
    style._element.rPr.rFonts.set(qn("w:eastAsia"), font_name)


def add_key_value_row(table: Any, key: str, value: str) -> None:
    cells = table.add_row().cells
    cells[0].text = key
    cells[1].text = value
    set_cell_shading(cells[0], "E9EEF5")


def build_offer_letter(
    *,
    company_name: str,
    hr_department: str,
    candidate_name: str,
    job_title: str,
    department: str,
    work_location: str,
    employment_type: str,
    expected_join_date: date,
    offer_date: date,
    acceptance_deadline: date,
    recognized_career_text: str,
    base_salary: float,
    performance_salary: float,
    fixed_overtime: float,
    incentive: float,
    cash_benefit_values: dict[str, float],
    sign_on_bonus: float,
    benefits: list[dict[str, Any]],
    special_terms: str,
) -> bytes:
    document = Document()
    set_korean_font(document)

    title = document.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title.add_run("OFFER LETTER")
    title_run.bold = True
    title_run.font.size = Pt(20)

    subtitle = document.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle_run = subtitle.add_run(f"{company_name} 채용 제안서")
    subtitle_run.bold = True
    subtitle_run.font.size = Pt(13)

    document.add_paragraph()
    intro = document.add_paragraph()
    intro.add_run(f"{candidate_name} 님께,").bold = True
    document.add_paragraph(
        f"당사는 아래와 같이 {candidate_name} 님께 입사를 제안드립니다. "
        "본 제안은 입사에 필요한 제반 절차와 최종 근로계약 체결을 전제로 합니다."
    )

    table = document.add_table(rows=1, cols=2)
    table.style = "Table Grid"
    table.rows[0].cells[0].text = "항목"
    table.rows[0].cells[1].text = "제안 내용"
    for cell in table.rows[0].cells:
        set_cell_shading(cell, "D9E2F3")

    add_key_value_row(table, "제안일", offer_date.strftime("%Y년 %m월 %d일"))
    add_key_value_row(table, "입사예정일", expected_join_date.strftime("%Y년 %m월 %d일"))
    add_key_value_row(table, "고용형태", employment_type)
    add_key_value_row(table, "부서", department or "추후 확정")
    add_key_value_row(table, "직무/직위", job_title or "추후 확정")
    add_key_value_row(table, "근무지", work_location or "추후 확정")
    add_key_value_row(table, "인정경력", recognized_career_text)

    document.add_paragraph()
    document.add_paragraph("1. 보상 조건").runs[0].bold = True

    compensation_table = document.add_table(rows=1, cols=3)
    compensation_table.style = "Table Grid"
    for index, value in enumerate(["구성항목", "연간 금액", "비고"]):
        compensation_table.rows[0].cells[index].text = value
        set_cell_shading(compensation_table.rows[0].cells[index], "D9E2F3")

    compensation_rows = [
        ("기본급", base_salary, "계약연봉 구성"),
        ("업적급", performance_salary, "계약연봉 구성"),
        ("고정연장수당", fixed_overtime, "계약연봉 구성"),
        ("성과급", incentive, "경영성과급 등, 변동 가능"),
    ]
    for label, amount, note in compensation_rows:
        cells = compensation_table.add_row().cells
        cells[0].text = label
        cells[1].text = format_won(amount)
        cells[2].text = note

    contract_salary = base_salary + performance_salary + fixed_overtime
    cash_total = sum(float(value or 0) for value in cash_benefit_values.values())
    annual_total = contract_salary + incentive + cash_total
    first_year_total = annual_total + sign_on_bonus

    cells = compensation_table.add_row().cells
    cells[0].text = "계약연봉"
    cells[1].text = format_won(contract_salary)
    cells[2].text = "기본급+업적급+고정연장수당"
    for cell in cells:
        set_cell_shading(cell, "E2F0D9")

    for label, amount in cash_benefit_values.items():
        if float(amount or 0) <= 0:
            continue
        cells = compensation_table.add_row().cells
        cells[0].text = label
        cells[1].text = format_won(amount)
        cells[2].text = "현금성지급 복리후생"

    cells = compensation_table.add_row().cells
    cells[0].text = "현금성지급 복리후생 소계"
    cells[1].text = format_won(cash_total)
    cells[2].text = "연간 기준"

    cells = compensation_table.add_row().cells
    cells[0].text = "총연봉"
    cells[1].text = format_won(annual_total)
    cells[2].text = "계약연봉+성과급+현금성지급 복리후생"
    for cell in cells:
        set_cell_shading(cell, "DDEBF7")

    if sign_on_bonus > 0:
        cells = compensation_table.add_row().cells
        cells[0].text = "입사 1차년도 총보상"
        cells[1].text = format_won(first_year_total)
        cells[2].text = "총연봉+사이닝 보너스"

    document.add_paragraph()
    document.add_paragraph("2. 주요 복리후생").runs[0].bold = True
    if benefits:
        for benefit in benefits:
            paragraph = document.add_paragraph(style="List Bullet")
            text = f"{benefit.get('복리후생', '')} ({benefit.get('적용여부', '')})"
            description = benefit.get("설명", benefit.get("비고", ""))
            if description:
                text += f" - {description}"
            paragraph.add_run(text)
    else:
        document.add_paragraph("복리후생은 회사 규정과 고용형태에 따라 적용됩니다.")

    document.add_paragraph()
    document.add_paragraph("3. 기타 조건").runs[0].bold = True
    document.add_paragraph(
        "성과급, 복리후생 및 기타 보상은 회사의 제도 변경, 지급기준, 재직요건, "
        "개인 및 조직 성과 등에 따라 달라질 수 있습니다."
    )
    document.add_paragraph(
        "본 제안서는 근로계약서가 아니며, 최종 근로조건은 입사 시 체결하는 "
        "근로계약서, 취업규칙 및 회사 제 규정에 따릅니다."
    )
    if special_terms.strip():
        document.add_paragraph(special_terms.strip())

    document.add_paragraph()
    document.add_paragraph(
        f"본 제안의 수락 여부를 {acceptance_deadline:%Y년 %m월 %d일}까지 "
        "회신하여 주시기 바랍니다."
    )

    closing = document.add_paragraph()
    closing.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    closing.add_run(f"{company_name}\n{hr_department}").bold = True

    document.add_paragraph()
    document.add_paragraph("[입사 제안 수락]").runs[0].bold = True
    document.add_paragraph(
        "본인은 위 제안 내용을 확인하였으며, 안내받은 조건과 절차에 따라 "
        "입사 제안을 수락합니다."
    )

    signature_table = document.add_table(rows=3, cols=2)
    signature_table.style = "Table Grid"
    values = [
        ("성명", candidate_name),
        ("서명", ""),
        ("일자", "        년      월      일"),
    ]
    for row_index, (label, value) in enumerate(values):
        signature_table.cell(row_index, 0).text = label
        signature_table.cell(row_index, 1).text = value
        set_cell_shading(signature_table.cell(row_index, 0), "E9EEF5")

    buffer = BytesIO()
    document.save(buffer)
    return buffer.getvalue()

# -----------------------------------------------------------------------------
# v3.7: 세로형 PPT 오퍼레터 생성
# -----------------------------------------------------------------------------
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt as PptPt


def _ppt_add_text(slide, text: str, x: float, y: float, w: float, h: float,
                  size: int = 14, bold: bool = False, color: str = "111111",
                  align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.name = "Malgun Gothic"
    run.font.color.rgb = RGBColor.from_string(color)
    return box


def _ppt_add_rect(slide, x: float, y: float, w: float, h: float, fill: str = "F4F7FB", line: str = "D9E2F3"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill)
    shape.line.color.rgb = RGBColor.from_string(line)
    return shape


def _ppt_add_kv(slide, label: str, value: str, x: float, y: float, w: float, h: float):
    _ppt_add_rect(slide, x, y, w, h, fill="FFFFFF", line="D9E2F3")
    _ppt_add_text(slide, label, x + 0.12, y + 0.08, w * 0.34, h - 0.1, size=10, bold=True, color="1F4E79")
    _ppt_add_text(slide, value, x + w * 0.37, y + 0.08, w * 0.58, h - 0.1, size=10, color="222222")


def build_offer_ppt(
    *,
    company_name: str,
    hr_department: str,
    candidate_name: str,
    job_title: str,
    department: str,
    work_location: str,
    employment_type: str,
    expected_join_date: date,
    offer_date: date,
    acceptance_deadline: date,
    recognized_career_text: str,
    base_salary: float,
    performance_salary: float,
    fixed_overtime: float,
    incentive: float,
    cash_benefit_values: dict[str, float],
    sign_on_bonus: float,
    benefits: list[dict[str, Any]],
    special_terms: str,
) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(7.5)
    prs.slide_height = Inches(13.333)
    blank = prs.slide_layouts[6]

    contract_salary = base_salary + performance_salary + fixed_overtime
    cash_total = sum(float(v or 0) for v in cash_benefit_values.values())
    annual_total = contract_salary + incentive + cash_total
    first_year_total = annual_total + sign_on_bonus

    # Slide 1: Summary
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor.from_string("F7F9FC")
    _ppt_add_text(slide, "OFFER LETTER", 0.55, 0.55, 6.4, 0.45, size=24, bold=True, color="1F4E79", align=PP_ALIGN.CENTER)
    _ppt_add_text(slide, f"{company_name} 채용 제안서", 0.55, 1.05, 6.4, 0.35, size=13, bold=True, color="333333", align=PP_ALIGN.CENTER)
    _ppt_add_rect(slide, 0.55, 1.75, 6.4, 2.3, fill="FFFFFF")
    _ppt_add_text(slide, f"{candidate_name} 님께", 0.85, 2.05, 5.8, 0.35, size=18, bold=True, color="111111")
    _ppt_add_text(slide, "아래와 같이 입사를 제안드립니다. 본 제안은 입사에 필요한 제반 절차와 최종 근로계약 체결을 전제로 합니다.", 0.85, 2.55, 5.8, 0.9, size=12, color="333333")
    y = 4.45
    pairs = [
        ("제안일", offer_date.strftime("%Y.%m.%d")),
        ("입사예정일", expected_join_date.strftime("%Y.%m.%d")),
        ("고용형태", employment_type),
        ("부서", department or "추후 확정"),
        ("직무/직위", job_title or "추후 확정"),
        ("근무지", work_location or "추후 확정"),
        ("인정경력", recognized_career_text),
    ]
    for label, value in pairs:
        _ppt_add_kv(slide, label, value, 0.65, y, 6.2, 0.62)
        y += 0.72
    _ppt_add_text(slide, f"{company_name} · {hr_department}", 0.65, 12.35, 6.2, 0.25, size=10, color="666666", align=PP_ALIGN.RIGHT)

    # Slide 2: Compensation
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
    _ppt_add_text(slide, "보상 조건", 0.55, 0.55, 6.4, 0.5, size=22, bold=True, color="1F4E79")
    comp_rows = [
        ("제안 기본급", base_salary),
        ("제안 업적급", performance_salary),
        ("제안 고정연장수당", fixed_overtime),
        ("계약연봉", contract_salary),
        ("성과급", incentive),
        ("현금성지급 복리후생", cash_total),
        ("총연봉", annual_total),
    ]
    if sign_on_bonus > 0:
        comp_rows.append(("입사 1차년도 총보상", first_year_total))
    y = 1.45
    for label, value in comp_rows:
        fill = "E2F0D9" if label in {"계약연봉", "총연봉"} else "F7F9FC"
        _ppt_add_rect(slide, 0.65, y, 6.2, 0.58, fill=fill)
        _ppt_add_text(slide, label, 0.85, y + 0.1, 3.1, 0.3, size=11, bold=label in {"계약연봉", "총연봉"}, color="111111")
        _ppt_add_text(slide, format_won(value), 4.0, y + 0.1, 2.65, 0.3, size=11, bold=True, color="1F4E79", align=PP_ALIGN.RIGHT)
        y += 0.68
    _ppt_add_text(slide, "※ 성과급, 복리후생 및 기타 보상은 회사의 제도 변경, 지급기준, 재직요건, 개인 및 조직 성과 등에 따라 달라질 수 있습니다.", 0.65, 11.65, 6.2, 0.6, size=9, color="666666")

    # Slide 3: Benefits
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor.from_string("F7F9FC")
    _ppt_add_text(slide, "주요 복리후생", 0.55, 0.55, 6.4, 0.5, size=22, bold=True, color="1F4E79")
    benefit_items = benefits[:12]
    y = 1.35
    if benefit_items:
        for item in benefit_items:
            title = str(item.get("복리후생", ""))
            desc = str(item.get("설명", item.get("비고", "")) or "")
            _ppt_add_rect(slide, 0.65, y, 6.2, 0.75, fill="FFFFFF")
            _ppt_add_text(slide, title, 0.85, y + 0.08, 5.8, 0.25, size=11, bold=True, color="111111")
            _ppt_add_text(slide, desc[:130], 0.85, y + 0.35, 5.8, 0.32, size=8, color="555555")
            y += 0.86
            if y > 11.2:
                break
    else:
        _ppt_add_text(slide, "복리후생은 회사 규정과 고용형태에 따라 적용됩니다.", 0.75, 1.4, 6.0, 0.5, size=12)

    # Slide 4: Acceptance
    slide = prs.slides.add_slide(blank)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
    _ppt_add_text(slide, "수락 안내", 0.55, 0.55, 6.4, 0.5, size=22, bold=True, color="1F4E79")
    terms = "본 제안서는 근로계약서가 아니며, 최종 근로조건은 입사 시 체결하는 근로계약서, 취업규칙 및 회사 제 규정에 따릅니다."
    if special_terms.strip():
        terms += "\n\n" + special_terms.strip()
    _ppt_add_rect(slide, 0.65, 1.45, 6.2, 3.2, fill="F7F9FC")
    _ppt_add_text(slide, terms, 0.9, 1.75, 5.7, 2.5, size=11, color="333333")
    _ppt_add_text(slide, f"본 제안의 수락 여부를 {acceptance_deadline:%Y.%m.%d}까지 회신하여 주시기 바랍니다.", 0.75, 5.2, 6.0, 0.5, size=13, bold=True, color="111111")
    _ppt_add_text(slide, "[입사 제안 수락]", 0.75, 6.3, 6.0, 0.35, size=14, bold=True)
    _ppt_add_text(slide, "본인은 위 제안 내용을 확인하였으며, 안내받은 조건과 절차에 따라 입사 제안을 수락합니다.", 0.75, 6.75, 6.0, 0.7, size=11)
    _ppt_add_kv(slide, "성명", candidate_name, 0.75, 8.1, 6.0, 0.65)
    _ppt_add_kv(slide, "서명", "", 0.75, 8.9, 6.0, 0.65)
    _ppt_add_kv(slide, "일자", "        년      월      일", 0.75, 9.7, 6.0, 0.65)
    _ppt_add_text(slide, f"{company_name}\n{hr_department}", 0.75, 11.5, 6.0, 0.8, size=12, bold=True, color="1F4E79", align=PP_ALIGN.RIGHT)

    out = BytesIO()
    prs.save(out)
    return out.getvalue()
